from pptx import Presentation
import os

slides_dir = '/home/wgy/reading_group/LLM-Reading-Group/Slides'
slides_list = os.listdir(slides_dir)
markdown_dir = '/home/wgy/reading_group/LLM-Reading-Group/Markdown'
markdown_list = os.listdir(markdown_dir)

for slides_name in slides_list:
    base_name = slides_name.split('.')[0]
    if base_name not in markdown_list:
        slides_path = os.path.join(slides_dir, slides_name)
        markdown_name = base_name + '.md'
        markdown_path = os.path.join(markdown_dir, markdown_name)

        images_path = os.path.join('/home/wgy/reading_group/LLM-Reading-Group/Markdown/Images', base_name)
        if not os.path.exists(images_path):
            os.makedirs(images_path)

        # 加载PPT文件
        presentation = Presentation(slides_path)

        # 创建一个列表来存储从PPT提取的文本和图片信息
        content = []

        # 读取PPT的每一页
        for slide_number, slide in enumerate(presentation.slides):
            index = 1
            # 读取每个形状
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # 形状是文本框，提取文本
                    content.append(shape.text)
                    content.append(('\n'))
                elif shape.shape_type == 13:
                    # 形状是图片，保存图片
                    image = shape.image
                    image_bytes = image.blob
                    image_path = os.path.join(images_path, f'img_{slide_number}_{index}.png')
                    index += 1
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                    # 将图片的Markdown路径写入列表
                    width_px = int(shape.width * 96 / 914400)
                    height_px = int(shape.height * 96 / 914400)
                    img_tag = f'<img src="{image_path}" width="{width_px}" height="{height_px}"/>'
                    content.append(img_tag)
                    content.append('\n')

        # 将提取的内容写入Markdown文件
        with open(markdown_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(content))

        print(base_name)
